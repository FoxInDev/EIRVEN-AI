from __future__ import annotations

import json
import mimetypes
import zipfile
from pathlib import Path
from typing import Iterable

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.gif', '.bmp'}
TEXT_EXTENSIONS = {
    '.txt', '.md', '.rst', '.log', '.csv', '.tsv', '.json', '.jsonl', '.xml', '.yaml', '.yml',
    '.toml', '.ini', '.cfg', '.conf', '.py', '.js', '.mjs', '.cjs', '.ts', '.tsx', '.jsx',
    '.html', '.htm', '.css', '.scss', '.less', '.java', '.kt', '.kts', '.c', '.h', '.cpp',
    '.hpp', '.cs', '.go', '.rs', '.php', '.rb', '.sh', '.ps1', '.bat', '.cmd', '.sql',
    '.dockerfile', '.gitignore', '.env', '.properties', '.gradle', '.swift', '.dart', '.lua',
}


def is_image(path: str | Path) -> bool:
    p = Path(path)
    media = mimetypes.guess_type(p.name)[0] or ''
    return media.startswith('image/') or p.suffix.casefold() in IMAGE_EXTENSIONS


def split_images(paths: Iterable[str]) -> tuple[list[str], list[str]]:
    images: list[str] = []
    others: list[str] = []
    for raw in paths:
        if is_image(raw):
            images.append(str(raw))
        else:
            others.append(str(raw))
    return images, others


def _read_text(path: Path, limit: int) -> str:
    data = path.read_bytes()[: max(limit * 4, limit)]
    for encoding in ('utf-8', 'utf-8-sig', 'cp1251', 'latin-1'):
        try:
            return data.decode(encoding)[:limit]
        except UnicodeDecodeError:
            continue
    return data.decode('utf-8', errors='replace')[:limit]


def _read_pdf(path: Path, limit: int) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    chunks: list[str] = []
    remaining = limit
    for index, page in enumerate(reader.pages[:24], start=1):
        if remaining <= 0:
            break
        text = (page.extract_text() or '').strip()
        if text:
            piece = f'--- страница {index} ---\n{text}'[:remaining]
            chunks.append(piece)
            remaining -= len(piece)
    return '\n\n'.join(chunks)


def _read_docx(path: Path, limit: int) -> str:
    from docx import Document  # type: ignore

    doc = Document(str(path))
    chunks: list[str] = []
    remaining = limit
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        piece = text[:remaining]
        chunks.append(piece)
        remaining -= len(piece)
        if remaining <= 0:
            break
    return '\n'.join(chunks)


def _read_xlsx(path: Path, limit: int) -> str:
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    remaining = limit
    try:
        for ws in wb.worksheets[:8]:
            header = f'--- лист: {ws.title} ---\n'
            chunks.append(header)
            remaining -= len(header)
            for row in ws.iter_rows(values_only=True):
                line = '\t'.join('' if value is None else str(value) for value in row).rstrip()
                if not line:
                    continue
                piece = (line + '\n')[:remaining]
                chunks.append(piece)
                remaining -= len(piece)
                if remaining <= 0:
                    break
            if remaining <= 0:
                break
    finally:
        wb.close()
    return ''.join(chunks)



def _read_pptx(path: Path, limit: int) -> str:
    from pptx import Presentation  # type: ignore

    deck = Presentation(str(path))
    chunks: list[str] = []
    remaining = limit
    for index, slide in enumerate(deck.slides[:40], start=1):
        if remaining <= 0:
            break
        lines: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                lines.append(text)
        if lines:
            piece = (f"--- слайд {index} ---\n" + "\n".join(lines) + "\n")[:remaining]
            chunks.append(piece)
            remaining -= len(piece)
    return "\n".join(chunks)


def _read_archive(path: Path, limit: int) -> str:
    if not zipfile.is_zipfile(path):
        return "Архив не является ZIP-совместимым; оригинал доступен инструментам по указанному пути."
    lines = ["Содержимое ZIP-совместимого архива:"]
    used = len(lines[0])
    candidates: list[zipfile.ZipInfo] = []
    with zipfile.ZipFile(path) as archive:
        infos = [i for i in archive.infolist()[:500] if not i.is_dir()]
        for info in infos[:260]:
            line = f"{info.filename}\t{info.file_size} байт"
            if used + len(line) + 1 > max(1200, limit // 3):
                break
            lines.append(line)
            used += len(line) + 1
            suffix = Path(info.filename).suffix.casefold()
            base = Path(info.filename).name.casefold()
            if (
                info.file_size <= 350_000
                and not any(part in {'.git', 'node_modules', '.venv', '__pycache__', 'dist', 'build'} for part in Path(info.filename).parts)
                and (suffix in TEXT_EXTENSIONS or base in {'dockerfile', 'makefile', 'readme', 'license'})
            ):
                candidates.append(info)
        # Prefer entry points/config/readmes and then source files. Never extract members
        # to disk: read bytes directly from the archive to avoid path traversal.
        priority = {'readme.md', 'pyproject.toml', 'package.json', 'requirements.txt', 'main.py', 'app.py', 'run.py', 'dockerfile'}
        candidates.sort(key=lambda i: (Path(i.filename).name.casefold() not in priority, len(Path(i.filename).parts), i.filename.casefold()))
        for info in candidates[:18]:
            if used >= limit:
                break
            budget = min(5000, limit - used)
            try:
                raw = archive.read(info)[: max(budget * 3, budget)]
                text = ''
                for encoding in ('utf-8', 'utf-8-sig', 'cp1251', 'latin-1'):
                    try:
                        text = raw.decode(encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                text = text[:budget].strip()
                if not text:
                    continue
                section = f"\n--- {info.filename} ---\n{text}"
                section = section[: limit-used]
                lines.append(section)
                used += len(section)
            except Exception as exc:
                note = f"\n--- {info.filename}: не удалось прочитать: {exc} ---"
                note = note[: max(0, limit-used)]
                lines.append(note); used += len(note)
    return "\n".join(lines)[:limit]

def extract_attachment_context(paths: Iterable[str], *, total_limit: int = 22_000) -> str:
    """Return a bounded prompt section for local attachments.

    Rich documents are parsed locally. Unknown/binary files are still exposed with their
    exact local path and metadata so EIRVEN's tool layer can open/process them when the
    user asks for an action involving that attachment.
    """
    chunks: list[str] = []
    remaining = total_limit
    for raw in list(paths)[:16]:
        if remaining <= 0:
            break
        path = Path(raw)
        if not path.is_file():
            continue
        if is_image(path):
            continue
        stat = path.stat()
        media = mimetypes.guess_type(path.name)[0] or 'application/octet-stream'
        header = f'\n--- вложение: {path.name} | {media} | {stat.st_size} байт | локальный путь: {path} ---\n'
        chunks.append(header[:remaining])
        remaining -= min(len(header), remaining)
        if remaining <= 0:
            break
        try:
            suffix = path.suffix.casefold()
            per_file = min(9_000, remaining)
            if suffix in TEXT_EXTENSIONS or media.startswith('text/'):
                content = _read_text(path, per_file)
            elif suffix == '.pdf':
                content = _read_pdf(path, per_file)
            elif suffix == '.docx':
                content = _read_docx(path, per_file)
            elif suffix in {'.xlsx', '.xlsm'}:
                content = _read_xlsx(path, per_file)
            elif suffix in {'.pptx', '.ppsx'}:
                content = _read_pptx(path, per_file)
            elif suffix in {'.zip', '.jar', '.whl', '.docx', '.xlsx', '.xlsm', '.pptx', '.ppsx'} and zipfile.is_zipfile(path):
                # OOXML is handled above. This branch mainly exposes ordinary ZIP/JAR/WHL
                # structure without extracting untrusted archive members.
                content = _read_archive(path, per_file)
            elif suffix in {'.json'}:
                content = json.dumps(json.loads(_read_text(path, per_file)), ensure_ascii=False, indent=2)[:per_file]
            else:
                content = (
                    'Бинарное/мультимедийное вложение сохранено локально. '
                    'Если требуется действие или более глубокий анализ, используй инструменты ОС/файлов по указанному пути.'
                )
        except Exception as exc:
            content = f'Локальный разбор не удался: {exc}. Файл всё равно доступен инструментам по указанному пути.'
        content = content[:remaining]
        chunks.append(content)
        remaining -= len(content)
    return ''.join(chunks).strip()
